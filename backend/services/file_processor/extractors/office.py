import io
import re
import zipfile
import xml.etree.ElementTree as ET


def extract(file_bytes: bytes, ext: str) -> str:
    # Modern OpenXML formats
    if ext == ".docx":
        return _extract_docx(file_bytes)
    if ext == ".pptx":
        return _extract_pptx(file_bytes)
    if ext == ".xlsx":
        return _extract_xlsx(file_bytes)
    # Legacy binary Office formats
    if ext == ".doc":
        return _extract_doc(file_bytes)
    if ext == ".ppt":
        return _extract_ppt(file_bytes)
    if ext == ".xls":
        return _extract_xls(file_bytes)
    # OpenDocument formats
    if ext in (".odt", ".odp", ".ods"):
        return _extract_opendocument(file_bytes)
    raise ValueError(f"Unsupported office format: {ext}")


# ── Modern OpenXML ────────────────────────────────────────────────────────────

def _extract_docx(file_bytes: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = [
            shape.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        if slide_texts:
            lines.append(f"[Slide {i}]")
            lines.extend(slide_texts)
    return "\n".join(lines)


def _extract_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(", ".join(cells))
    wb.close()
    return "\n".join(lines)


# ── Legacy Binary Office ──────────────────────────────────────────────────────

def _extract_doc(file_bytes: bytes) -> str:
    """Best-effort text extraction from old .doc (OLE compound document)."""
    try:
        import olefile

        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        if ole.exists("WordDocument"):
            raw = ole.openstream("WordDocument").read()
            # Decode as UTF-16-LE and strip non-printable characters
            text = raw.decode("utf-16-le", errors="ignore")
            text = re.sub(r"[^\x20-\x7E\n\r\t\u00A0-\uFFFF]", " ", text)
            text = re.sub(r" {3,}", " ", text).strip()
        else:
            text = ""
        ole.close()
        return text if text else "Legacy Word document (.doc) — text extraction limited"
    except Exception as e:
        return f"Legacy Word document (.doc) — extraction failed: {e}"


def _extract_ppt(file_bytes: bytes) -> str:
    """Best-effort text extraction from old .ppt (OLE compound document)."""
    try:
        import olefile

        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        texts: list[str] = []
        for entry in ole.listdir():
            stream_name = "/".join(entry)
            if "PowerPoint Document" in stream_name or "Current User" in stream_name:
                try:
                    raw = ole.openstream(entry).read()
                    text = raw.decode("utf-16-le", errors="ignore")
                    text = re.sub(r"[^\x20-\x7E\n\r\t]", " ", text)
                    text = re.sub(r" {3,}", " ", text).strip()
                    if text:
                        texts.append(text)
                except Exception:
                    pass
        ole.close()
        combined = "\n".join(texts).strip()
        return combined if combined else "Legacy PowerPoint presentation (.ppt) — text extraction limited"
    except Exception as e:
        return f"Legacy PowerPoint presentation (.ppt) — extraction failed: {e}"


def _extract_xls(file_bytes: bytes) -> str:
    """Extract cell data from old .xls using xlrd."""
    try:
        import xlrd

        wb = xlrd.open_workbook(file_contents=file_bytes)
        lines: list[str] = []
        for sheet in wb.sheets():
            lines.append(f"[Sheet: {sheet.name}]")
            for row_idx in range(sheet.nrows):
                cells = [
                    str(sheet.cell_value(row_idx, col))
                    for col in range(sheet.ncols)
                    if sheet.cell_value(row_idx, col) != ""
                ]
                if cells:
                    lines.append(", ".join(cells))
        return "\n".join(lines)
    except Exception as e:
        return f"Legacy Excel workbook (.xls) — extraction failed: {e}"


# ── OpenDocument ──────────────────────────────────────────────────────────────

# OpenDocument XML namespaces
_OD_NS = {
    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
}


def _extract_opendocument(file_bytes: bytes) -> str:
    """Extract text from .odt / .odp / .ods (OpenDocument ZIP + XML)."""
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            content_xml = zf.read("content.xml")

        root = ET.fromstring(content_xml)
        texts: list[str] = []

        # Text paragraphs (odt, odp)
        for p in root.iter(f"{{{_OD_NS['text']}}}p"):
            text = "".join(p.itertext()).strip()
            if text:
                texts.append(text)

        # Table cells (ods)
        for cell in root.iter(f"{{{_OD_NS['table']}}}table-cell"):
            for p in cell.iter(f"{{{_OD_NS['text']}}}p"):
                text = "".join(p.itertext()).strip()
                if text:
                    texts.append(text)

        return "\n".join(texts) if texts else "OpenDocument file — no text content found"
    except Exception as e:
        return f"OpenDocument file — extraction failed: {e}"
